# import pptx  # pip install python-pptx
# from pptx.util import Inches
# from pptxtopdf import convert  # pip install pptxtopdf
#
#
# def input_data_and_save_pdf(template_path, output_path, data):
#     """
#     Inputs data into placeholders in a PowerPoint template and saves it as a PDF.
#
#     Args:
#         template_path: Path to the PowerPoint template file.
#         output_path: Path to save the output PDF file.
#         data: Dictionary containing data to be inserted into placeholders.
#     """
#
#     # Open the template presentation
#     prs = pptx.Presentation(template_path)
#
#     # Iterate through slides and find placeholders
#     for slide in prs.slides:
#         for shape in slide.shapes:
#             if shape.has_text_frame:
#                 text_frame = shape.text_frame
#                 for paragraph in text_frame.paragraphs:
#                     for run in paragraph.runs:
#                         placeholder_text = run.text
#                         if placeholder_text in data:
#                             run.text = data[placeholder_text]
#
#     # Save the presentation as a PDF
#     prs.save(output_path)
#
#
# # Example usage
# template_path = "input_files/File Format for Control Union.pptx"
# output_path = "output_files/CU_report_from_pptx.pptx"
# data = {
#     "deforestation_text": "Deforestation risk is low",
#     "encroachment_text1": "Encroachment risk is low ",
#     "deforestation_text2": "Deforestation risk is low",
#     "encroachment_text2": "Encroachment risk is low ",
#     "deforestation_text3": "Deforestation risk is low",
#     "encroachment_text3": "Encroachment risk is low ",
#     "deforestation_text4": "Deforestation risk is low",
#     "encroachment_text4": "Encroachment risk is low ",
#     "total_area_val": "0.12 ha",
#     "potec_val": "0.00 ha",
#     "def_val": "0.00ha",
#     "eligible_area_val": "0.12ha",
#     "tec_val": "20M"
# }
#
# input_data_and_save_pdf(template_path, output_path, data)
#
#
# input_dir = output_path
# output_dir = r"output_files"
#
# convert(input_dir, output_dir)


###################################################################################################


# import pptx  # pip install python-pptx
# from pptx.util import Inches
# from pptxtopdf import convert  # pip install pptxtopdf
#
#
# def input_data_and_save_pdf(template_path, output_path, data):
#     """
#     Inputs data into placeholders in a PowerPoint template and saves it as a PDF.
#
#     Args:
#         template_path: Path to the PowerPoint template file.
#         output_path: Path to save the output PDF file.
#         data: Dictionary containing data to be inserted into placeholders.
#     """
#
#     # Open the template presentation
#     prs = pptx.Presentation(template_path)
#
#     # Iterate through slides and find placeholders
#     for slide in prs.slides:
#         for shape in slide.shapes:
#             if shape.has_text_frame:
#                 text_frame = shape.text_frame
#                 for paragraph in text_frame.paragraphs:
#                     for run in paragraph.runs:
#                         placeholder_text = run.text
#                         print(f"Placeholder text found: {placeholder_text}")  # Debug print
#                         if placeholder_text in data:
#                             run.text = data[placeholder_text]
#                             print(f"Replaced with: {run.text}")  # Debug print
#
#     # Save the presentation as a PowerPoint file
#     prs.save(output_path)
#
#
# # Example usage
# template_path = "input_files/File Format for Control Union.pptx"
# output_path = "output_files/CU_report_from_pptx.pptx"
# data = {
#     "deforestation_text": "Deforestation risk is lowwwww",
#     "encroachment_text1": "Encroachment risk is veryyyyy low",
#     "deforestation_text2": "Deforestation risk is low",
#     "encroachment_text2": "Encroachment risk is lowwwww",
#     "deforestation_text3": "Deforestation risk is low",
#     "encroachment_text3": "Encroachment risk is low",
#     "deforestation_text4": "Deforestation risk is low",
#     "encroachment_text4": "Encroachment risk is low",
#     "total_area_val": "0.12 ha",
#     "potec_val": "0.00 ha",
#     "def_val": "0.00ha",
#     "eligible_area_val": "0.12ha",
#     "tec_val": "20M"
# }
#
# input_data_and_save_pdf(template_path, output_path, data)
#
# # Convert PowerPoint to PDF
# input_dir = output_path
# output_dir = r"output_files"
#
# convert(input_dir, output_dir)



##########################################################################################################

import pptx  # pip install python-pptx
from pptx.util import Inches
from pptxtopdf import convert  # pip install pptxtopdf


def input_data_and_save_pdf(template_path, output_path, data, image_data):
    """
    Inputs data into placeholders in a PowerPoint template and saves it as a PDF.
    Also, inserts images into specified shapes with given text.

    Args:
        template_path: Path to the PowerPoint template file.
        output_path: Path to save the output PowerPoint file.
        data: Dictionary containing text data to be inserted into placeholders.
        image_data: Dictionary containing image paths where the key is the shape text and the value is the image path.
    """

    # Open the template presentation
    prs = pptx.Presentation(template_path)

    # Iterate through slides and find placeholders
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        placeholder_text = run.text
                        if placeholder_text in data:
                            run.text = data[placeholder_text]
                        # Check for the specific text "image_1" and replace it with an image
                        elif placeholder_text in image_data:
                            print(f"Found placeholder for image: {placeholder_text}")
                            image_path = image_data[placeholder_text]
                            # Clear the text from the shape
                            run.text = ''
                            # Insert the image using the placeholder shape dimensions
                            left = shape.left
                            top = shape.top
                            width = shape.width
                            height = shape.height
                            print(f"Inserting image: {image_path} at {left}, {top}")
                            # Add the image with the same dimensions as the shape
                            slide.shapes.add_picture(image_path, left, top, width, height)

    # Save the presentation as a PowerPoint file
    prs.save(output_path)


# Example usage
template_path = "input_files/File Format for Control Union.pptx"
output_path = "output_files/CU_report_from_pptx.pptx"

# Text placeholders data
data = {
    "deforestation_text": "Deforestation risk is lowwwww",
    "encroachment_text1": "Encroachment risk is veryyyyy low",
    "deforestation_text2": "Deforestation risk is low",
    "encroachment_text2": "Encroachment risk is lowwwww",
    "deforestation_text3": "Deforestation risk is low",
    "encroachment_text3": "Encroachment risk is low",
    "deforestation_text4": "Deforestation risk is low",
    "encroachment_text4": "Encroachment risk is low",
    "total_area_val": "0.12 ha",
    "potec_val": "0.00 ha",
    "def_val": "0.00ha",
    "eligible_area_val": "0.12ha",
    "tec_val": "20M"
}

# Image placeholder data: key is the placeholder text (e.g., "image_1"), and the value is the image path.
image_data = {
    "image_1": "imgs/Picture1.png",
    "image_2": "imgs/Picture2.png",
    "image_3": "imgs/Picture3.png",
    "image_4": "imgs/Picture4.png"
}

# Call the function
input_data_and_save_pdf(template_path, output_path, data, image_data)

# Convert PowerPoint to PDF
input_dir = output_path
output_dir = r"output_files"

convert(input_dir, output_dir)
